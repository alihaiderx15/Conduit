
from pathlib import Path
from types import SimpleNamespace
import json
import pytest
from conduit.conversation.session import ConversationSession
from conduit.file_processing import FileProcessingService

def make_service(tmp_path):
    return FileProcessingService(state_path=tmp_path/"state.json")

def make_session():
    s=object.__new__(ConversationSession)
    s.agent=SimpleNamespace(loop=SimpleNamespace(provider=None,model="none"))
    s._file_context={}
    return s

@pytest.mark.asyncio
async def test_spreadsheet_analyze_chat(tmp_path,monkeypatch):
    import pandas as pd
    from conduit.conversation import session as sm
    p=tmp_path/"sales.xlsx"
    pd.DataFrame({"Revenue":[100,200],"Region":["N","S"]}).to_excel(p,index=False)
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    async def fake(r): return "Revenue varies across the sample and the spreadsheet has no missing values."
    monkeypatch.setattr(s,"_complete_file_semantic_result",fake)
    assert ConversationSession._could_be_file_processing_request("Analyze")
    answer,report=await s._execute_file_processing_request("Analyze")
    assert report.success and "Revenue varies" in answer
    assert not list(tmp_path.glob("*_analysis*.txt"))

@pytest.mark.asyncio
async def test_sort_missing_column_and_followup(tmp_path,monkeypatch):
    import pandas as pd
    from conduit.conversation import session as sm
    p=tmp_path/"sales.xlsx"
    pd.DataFrame({"Revenue":[100,300],"Region":["N","S"]}).to_excel(p,index=False)
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    answer,report=await s._execute_file_processing_request("Sort this file")
    assert not report.success and "Which column" in answer
    answer,report=await s._continue_pending_file_operation("Revenue descending")
    assert report.success and "Sorted spreadsheet by Revenue" in answer
    out=next(tmp_path.glob("sales_sorted*.xlsx"))
    assert pd.read_excel(out)["Revenue"].tolist()==[300,100]

@pytest.mark.asyncio
async def test_filter_missing_condition_asks(tmp_path,monkeypatch):
    import pandas as pd
    from conduit.conversation import session as sm
    p=tmp_path/"sales.xlsx"; pd.DataFrame({"Revenue":[100,300]}).to_excel(p,index=False)
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    answer,report=await s._execute_file_processing_request("Filter this file")
    assert not report.success and "What filter should I apply" in answer

@pytest.mark.asyncio
async def test_spreadsheet_validate_clean_message(tmp_path,monkeypatch):
    import pandas as pd
    from conduit.conversation import session as sm
    p=tmp_path/"sales.xlsx"; pd.DataFrame({"A":[1]}).to_excel(p,index=False)
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    answer,report=await s._execute_file_processing_request("Validate this file")
    assert not report.success and "Validation isn't supported for spreadsheet files" in answer

@pytest.mark.asyncio
async def test_json_validate_format_convert_csv(tmp_path,monkeypatch):
    from conduit.conversation import session as sm
    p=tmp_path/"data.json"; p.write_text(json.dumps([{"a":1},{"a":2}]),encoding="utf-8")
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    answer,report=await s._execute_file_processing_request("Validate this file")
    assert report.success and answer=="JSON is valid."
    answer,report=await s._execute_file_processing_request("Format this file")
    assert report.success and "Saved output to" in answer
    answer,report=await s._execute_file_processing_request("convert this to csv")
    assert report.success and "Converted JSON to CSV" in answer

@pytest.mark.asyncio
async def test_json_analyze_chat(tmp_path,monkeypatch):
    from conduit.conversation import session as sm
    p=tmp_path/"data.json"; p.write_text('{"enabled":true}',encoding="utf-8")
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    async def fake(r): return "The JSON contains one enabled boolean setting."
    monkeypatch.setattr(s,"_complete_file_semantic_result",fake)
    answer,report=await s._execute_file_processing_request("Analyze this file")
    assert report.success and answer.startswith("The JSON")
    assert not list(tmp_path.glob("*_analysis*.txt"))

@pytest.mark.asyncio
async def test_ppt_extract_and_analyze(tmp_path,monkeypatch):
    from pptx import Presentation
    from conduit.conversation import session as sm
    p=tmp_path/"deck.pptx"
    prs=Presentation(); slide=prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text="Testing Conduit"; slide.placeholders[1].text="Extraction and analysis"; prs.save(p)
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    answer,report=await s._execute_file_processing_request("Extract text from this ppt")
    assert report.success and "Testing Conduit" in answer
    async def fake(r): return "The presentation is a short test deck for extraction and analysis."
    monkeypatch.setattr(s,"_complete_file_semantic_result",fake)
    answer,report=await s._execute_file_processing_request("Analyze this ppt")
    assert report.success and "short test deck" in answer

@pytest.mark.asyncio
async def test_image_describe_file_route(tmp_path,monkeypatch):
    from PIL import Image
    from conduit.conversation import session as sm
    p=tmp_path/"photo.jpg"; Image.new("RGB",(10,10)).save(p)
    service=make_service(tmp_path); service.register_dropped_file(p)
    monkeypatch.setattr(sm,"file_service",service)
    s=make_session()
    async def fake(r): return "The image is a plain test square."
    monkeypatch.setattr(s,"_complete_file_semantic_result",fake)
    assert ConversationSession._could_be_file_processing_request("describe this image")
    answer,report=await s._execute_file_processing_request("describe this image")
    assert report.success and "plain test square" in answer

def test_version():
    root=Path(__file__).resolve().parents[1]
    assert 'version = "3.1.8"' in (root/"pyproject.toml").read_text(encoding="utf-8")
