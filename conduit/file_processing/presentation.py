
from __future__ import annotations

from pathlib import Path
from typing import Any

from .common import DependencyUnavailable, FileProcessingError, file_basic_info, safe_output_path
from .models import FileInput, ProcessingResult


def extract_text(path: Path) -> tuple[str, int]:
    if path.suffix.casefold() != ".pptx":
        raise FileProcessingError("Presentation processing currently supports PPTX.")
    try:
        from pptx import Presentation
    except Exception as exc:
        raise DependencyUnavailable("Presentation processing requires python-pptx.") from exc

    prs = Presentation(str(path))
    blocks = []
    for index, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        blocks.append(f"[Slide {index}]\n" + "\n".join(slide_text))
    return "\n\n".join(blocks), len(prs.slides)


def process(file: FileInput, action: str, params: dict[str, Any]) -> ProcessingResult:
    text, slides = extract_text(file.path)
    action = action.casefold().strip()

    if action == "inspect":
        data = file_basic_info(file.path)
        data.update({"slides": slides, "characters": len(text)})
        return ProcessingResult(True, action, f"Inspected presentation with {slides} slide(s).", file, data=data)

    if action == "extract_text":
        target = safe_output_path(file.path, "extracted_text", ".txt")
        target.write_text(text, encoding="utf-8")
        return ProcessingResult(True, action, f"Extracted text from {slides} slide(s).", file,
                                output_path=target, data={"text": text, "slides": slides})

    if action in {"summarize", "analyze"}:
        instruction = str(params.get("instruction") or (
            "Summarize this presentation slide by slide, then give the overall key points."
            if action == "summarize"
            else "Analyze this presentation's structure, message, content quality, and key findings."
        ))
        return ProcessingResult(True, action, f"Prepared {slides} slide(s) for {action}.", file,
                                data={"slides": slides}, semantic_text=text,
                                semantic_instruction=instruction)

    raise FileProcessingError(f"Unsupported presentation action: {action}")
